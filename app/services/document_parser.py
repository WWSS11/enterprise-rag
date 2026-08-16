from __future__ import annotations

import csv
import json
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import xlrd
from bs4 import BeautifulSoup, Tag
from defusedxml import ElementTree
from docx import Document as WordDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

SUPPORTED_EXTENSIONS = frozenset(
    {
        ".txt",
        ".md",
        ".csv",
        ".json",
        ".xml",
        ".pdf",
        ".docx",
        ".pptx",
        ".xlsx",
        ".xls",
        ".html",
        ".htm",
    }
)

MARKDOWN_HEADING = re.compile(r"^(#{1,6})\s+(.+?)\s*#*\s*$")
WORD_HEADING = re.compile(r"^Heading\s+([1-6])$", re.IGNORECASE)


class UnsupportedDocumentError(ValueError):
    pass


@dataclass(frozen=True, slots=True)
class ParsedSection:
    text: str
    title: str | None = None
    heading_path: tuple[str, ...] = ()
    section_type: str = "prose"
    metadata: dict[str, Any] = field(default_factory=dict)


def _clean(text: str) -> str:
    text = text.replace("\x00", "")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _flat_section(text: str, *, section_type: str = "prose") -> list[ParsedSection]:
    cleaned = _clean(text)
    if not cleaned:
        return []
    metadata = {"paragraph_start": 1, "paragraph_end": 1} if section_type == "prose" else {}
    return [ParsedSection(text=cleaned, section_type=section_type, metadata=metadata)]


def _excel_column(column: int) -> str:
    output = ""
    while column > 0:
        column, remainder = divmod(column - 1, 26)
        output = chr(65 + remainder) + output
    return output or "A"


def _cell_range(first_column: int, last_column: int, row: int) -> str:
    start = f"{_excel_column(first_column)}{row}"
    end = f"{_excel_column(last_column)}{row}"
    return start if start == end else f"{start}:{end}"


def _parse_plain_text(path: Path) -> list[ParsedSection]:
    paragraphs = [part for part in re.split(r"\n\s*\n", _read_text(path)) if _clean(part)]
    return [
        ParsedSection(
            text=_clean(paragraph),
            metadata={"paragraph_start": index, "paragraph_end": index},
        )
        for index, paragraph in enumerate(paragraphs, start=1)
    ]


def _parse_csv(path: Path) -> list[ParsedSection]:
    sections: list[ParsedSection] = []
    with path.open("r", encoding="utf-8", errors="replace", newline="") as stream:
        for row_number, row in enumerate(csv.reader(stream), start=1):
            populated = [index for index, value in enumerate(row, start=1) if _clean(value)]
            if not populated:
                continue
            first_column, last_column = min(populated), max(populated)
            text = "\t".join(row[first_column - 1 : last_column]).rstrip()
            sections.append(
                ParsedSection(
                    text=text,
                    title=f"第 {row_number} 行",
                    section_type="table",
                    metadata={
                        "row_start": row_number,
                        "row_end": row_number,
                        "cell_range": _cell_range(first_column, last_column, row_number),
                    },
                )
            )
    return sections


def _parse_markdown(path: Path) -> list[ParsedSection]:
    sections: list[ParsedSection] = []
    heading_stack: list[str] = []
    title: str | None = None
    body: list[str] = []
    in_code_fence = False
    paragraph_number = 0

    def flush() -> None:
        nonlocal body, paragraph_number
        text = _clean("\n".join(body))
        for paragraph in (part for part in re.split(r"\n\s*\n", text) if _clean(part)):
            paragraph_number += 1
            sections.append(
                ParsedSection(
                    text=_clean(paragraph),
                    title=title,
                    heading_path=tuple(heading_stack),
                    section_type="prose",
                    metadata={
                        "paragraph_start": paragraph_number,
                        "paragraph_end": paragraph_number,
                    },
                )
            )
        body = []

    for line in _read_text(path).splitlines():
        if line.lstrip().startswith("```") or line.lstrip().startswith("~~~"):
            in_code_fence = not in_code_fence
            body.append(line)
            continue
        match = None if in_code_fence else MARKDOWN_HEADING.match(line)
        if not match:
            body.append(line)
            continue
        flush()
        level = len(match.group(1))
        heading = _clean(match.group(2))
        heading_stack = heading_stack[: level - 1]
        heading_stack.append(heading)
        title = heading
    flush()
    return sections


def _parse_json(path: Path) -> list[ParsedSection]:
    try:
        value = json.loads(_read_text(path))
    except json.JSONDecodeError:
        return _flat_section(_read_text(path), section_type="structured")
    pretty = json.dumps(value, ensure_ascii=False, indent=2)
    return _flat_section(pretty, section_type="structured")


def _parse_pdf(path: Path) -> list[ParsedSection]:
    reader = PdfReader(str(path))
    sections: list[ParsedSection] = []
    for page_number, page in enumerate(reader.pages, start=1):
        text = _clean(page.extract_text() or "")
        if text:
            sections.append(
                ParsedSection(
                    text=text,
                    title=f"第 {page_number} 页",
                    heading_path=(f"第 {page_number} 页",),
                    metadata={"page": page_number},
                )
            )
    return sections


def _parse_docx(path: Path) -> list[ParsedSection]:
    document = WordDocument(str(path))
    sections: list[ParsedSection] = []
    heading_stack: list[str] = []
    title: str | None = None
    for paragraph_number, paragraph in enumerate(document.paragraphs, start=1):
        text = _clean(paragraph.text)
        if not text:
            continue
        style_name = getattr(paragraph.style, "name", "") or ""
        match = WORD_HEADING.match(style_name)
        if match:
            level = int(match.group(1))
            heading_stack = heading_stack[: level - 1]
            heading_stack.append(text)
            title = text
        else:
            sections.append(
                ParsedSection(
                    text=text,
                    title=title,
                    heading_path=tuple(heading_stack),
                    metadata={
                        "paragraph_start": paragraph_number,
                        "paragraph_end": paragraph_number,
                    },
                )
            )

    for table_number, table in enumerate(document.tables, start=1):
        for row_number, row in enumerate(table.rows, start=1):
            values = [_clean(cell.text) for cell in row.cells]
            populated = [index for index, value in enumerate(values, start=1) if value]
            if not populated:
                continue
            first_column, last_column = min(populated), max(populated)
            text = "\t".join(values[first_column - 1 : last_column]).rstrip()
            table_title = f"表格 {table_number}"
            sections.append(
                ParsedSection(
                    text=text,
                    title=table_title,
                    heading_path=(*heading_stack, table_title),
                    section_type="table",
                    metadata={
                        "table": table_title,
                        "row_start": row_number,
                        "row_end": row_number,
                        "cell_range": _cell_range(
                            first_column, last_column, row_number
                        ),
                    },
                )
            )
    return sections


def _parse_pptx(path: Path) -> list[ParsedSection]:
    presentation = Presentation(str(path))
    sections: list[ParsedSection] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        title = _clean(slide.shapes.title.text) if slide.shapes.title else ""
        lines: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = _clean(shape.text)
                if text and text != title:
                    lines.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    values = [_clean(cell.text) for cell in row.cells]
                    if any(values):
                        lines.append("\t".join(values))
        text = _clean("\n".join(([title] if title else []) + lines))
        if text or title:
            display_title = title or f"幻灯片 {slide_number}"
            sections.append(
                ParsedSection(
                    text=text or display_title,
                    title=display_title,
                    heading_path=(display_title,),
                    metadata={"slide": slide_number},
                )
            )
    return sections


def _parse_xlsx(path: Path) -> list[ParsedSection]:
    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        sections: list[ParsedSection] = []
        for sheet in workbook.worksheets:
            for row_number, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                values = ["" if cell is None else str(cell) for cell in row]
                populated = [index for index, value in enumerate(values, start=1) if value]
                if not populated:
                    continue
                first_column, last_column = min(populated), max(populated)
                text = "\t".join(values[first_column - 1 : last_column]).rstrip()
                sections.append(
                    ParsedSection(
                        text=text,
                        title=sheet.title,
                        heading_path=(sheet.title,),
                        section_type="table",
                        metadata={
                            "sheet": sheet.title,
                            "row_start": row_number,
                            "row_end": row_number,
                            "cell_range": _cell_range(
                                first_column, last_column, row_number
                            ),
                        },
                    )
                )
        return sections
    finally:
        workbook.close()


def _parse_xls(path: Path) -> list[ParsedSection]:
    workbook = xlrd.open_workbook(str(path), on_demand=True)
    try:
        sections: list[ParsedSection] = []
        for sheet in workbook.sheets():
            for row_index in range(sheet.nrows):
                values = [
                    str(sheet.cell_value(row_index, column))
                    for column in range(sheet.ncols)
                ]
                populated = [index for index, value in enumerate(values, start=1) if value]
                if not populated:
                    continue
                first_column, last_column = min(populated), max(populated)
                row_number = row_index + 1
                text = "\t".join(values[first_column - 1 : last_column]).rstrip()
                sections.append(
                    ParsedSection(
                        text=text,
                        title=sheet.name,
                        heading_path=(sheet.name,),
                        section_type="table",
                        metadata={
                            "sheet": sheet.name,
                            "row_start": row_number,
                            "row_end": row_number,
                            "cell_range": _cell_range(
                                first_column, last_column, row_number
                            ),
                        },
                    )
                )
        return sections
    finally:
        workbook.release_resources()


def _parse_html(path: Path) -> list[ParsedSection]:
    soup = BeautifulSoup(_read_text(path), "html.parser")
    page_title = _clean(soup.title.get_text(" ", strip=True)) if soup.title else ""
    overview_parts: list[str] = []
    if page_title:
        overview_parts.append(f"页面标题：{page_title}")
    brands = list(
        dict.fromkeys(
            text
            for element in soup.select('[class*="brand"]')
            if (text := _clean(element.get_text(" ", strip=True)))
        )
    )
    if brands:
        overview_parts.append(f"页面品牌：{'；'.join(brands[:5])}")
    major_areas = list(
        dict.fromkeys(
            text
            for element in soup.find_all(["h1", "h2", "h3"])
            if (text := _clean(element.get_text(" ", strip=True)))
        )
    )
    if major_areas:
        overview_parts.append(f"主要界面区域：{'；'.join(major_areas[:20])}")

    for tag in soup(["script", "style", "noscript", "nav", "footer", "header"]):
        tag.decompose()

    sections: list[ParsedSection] = []
    if overview_parts:
        sections.append(
            ParsedSection(
                text="\n".join(overview_parts),
                title=page_title or "页面概览",
                heading_path=((page_title or "页面概览"),),
                section_type="metadata",
                metadata={"html_title": page_title},
            )
        )
    headings: list[str] = []
    current_title: str | None = None
    body: list[str] = []

    def flush() -> None:
        nonlocal body
        text = _clean("\n".join(body))
        if text:
            sections.append(
                ParsedSection(
                    text=text,
                    title=current_title,
                    heading_path=tuple(headings),
                )
            )
        body = []

    root = soup.body or soup
    for node in root.descendants:
        if not isinstance(node, Tag):
            continue
        if node.name in {"h1", "h2", "h3", "h4", "h5", "h6"}:
            flush()
            level = int(node.name[1])
            current_title = _clean(node.get_text(" ", strip=True))
            headings = headings[: level - 1]
            headings.append(current_title)
        elif node.name in {"p", "li", "pre", "blockquote"}:
            text = _clean(node.get_text(" ", strip=True))
            if text:
                body.append(text)
        elif node.name == "table":
            rows = [
                "\t".join(
                    _clean(cell.get_text(" ", strip=True))
                    for cell in row.find_all(["th", "td"])
                )
                for row in node.find_all("tr")
            ]
            table_text = _clean("\n".join(row for row in rows if row))
            if table_text:
                body.append(table_text)
    flush()
    if sections:
        return sections
    return _flat_section(soup.get_text("\n", strip=True))


def _parse_xml(path: Path) -> list[ParsedSection]:
    root = ElementTree.parse(path).getroot()
    text = _clean("\n".join(part.strip() for part in root.itertext() if part.strip()))
    return _flat_section(text, section_type="structured")


def parse_document_sections(path: Path) -> list[ParsedSection]:
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise UnsupportedDocumentError(f"unsupported document type: {suffix or '<none>'}")
    if suffix == ".md":
        return _parse_markdown(path)
    if suffix == ".txt":
        return _parse_plain_text(path)
    if suffix == ".csv":
        return _parse_csv(path)
    if suffix == ".json":
        return _parse_json(path)
    if suffix == ".xml":
        return _parse_xml(path)
    if suffix == ".pdf":
        return _parse_pdf(path)
    if suffix == ".docx":
        return _parse_docx(path)
    if suffix == ".pptx":
        return _parse_pptx(path)
    if suffix == ".xlsx":
        return _parse_xlsx(path)
    if suffix == ".xls":
        return _parse_xls(path)
    return _parse_html(path)


def parse_document(path: Path) -> str:
    """Compatibility helper returning the complete plain text representation."""

    return "\n\n".join(section.text for section in parse_document_sections(path))
