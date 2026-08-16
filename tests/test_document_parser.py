from app.services.document_parser import parse_document, parse_document_sections


def test_parse_text_document(tmp_path) -> None:
    path = tmp_path / "notes.txt"
    path.write_text("企业级 RAG", encoding="utf-8")
    assert parse_document(path) == "企业级 RAG"


def test_parse_xml_document(tmp_path) -> None:
    path = tmp_path / "knowledge.xml"
    path.write_text(
        "<root><title>权限设计</title><body>最小权限原则</body></root>",
        encoding="utf-8",
    )
    assert parse_document(path) == "权限设计\n最小权限原则"


def test_parse_powerpoint_document(tmp_path) -> None:
    from pptx import Presentation

    path = tmp_path / "architecture.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "企业 RAG"
    slide.placeholders[1].text = "异步入库与权限过滤"
    presentation.save(path)

    content = parse_document(path)
    assert "企业 RAG" in content
    assert "异步入库与权限过滤" in content


def test_parse_markdown_preserves_heading_hierarchy(tmp_path) -> None:
    path = tmp_path / "architecture.md"
    path.write_text(
        "# 企业知识库\n\n总体说明。\n\n## 检索层\n\n混合召回与重排。",
        encoding="utf-8",
    )

    sections = parse_document_sections(path)

    assert [section.heading_path for section in sections] == [
        ("企业知识库",),
        ("企业知识库", "检索层"),
    ]
    assert sections[1].title == "检索层"


def test_parse_html_preserves_page_title_and_major_areas(tmp_path) -> None:
    path = tmp_path / "legacy.html"
    path.write_text(
        """
        <html><head><title>RAG 智能问答</title></head>
        <body><a class="sidebar-brand"><span>RAG Chat</span></a>
        <h2>知识库</h2><p>上传并管理文档</p></body></html>
        """,
        encoding="utf-8",
    )

    sections = parse_document_sections(path)

    assert "页面标题：RAG 智能问答" in sections[0].text
    assert "页面品牌：RAG Chat" in sections[0].text
    assert "主要界面区域：知识库" in sections[0].text


def test_parse_docx_preserves_paragraph_and_table_cell_locations(tmp_path) -> None:
    from docx import Document

    path = tmp_path / "policy.docx"
    document = Document()
    document.add_heading("访问控制", level=1)
    document.add_paragraph("第一段规则。")
    document.add_paragraph("第二段规则。")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "角色"
    table.cell(0, 1).text = "权限"
    document.save(path)

    sections = parse_document_sections(path)

    assert sections[0].metadata == {"paragraph_start": 2, "paragraph_end": 2}
    assert sections[1].metadata == {"paragraph_start": 3, "paragraph_end": 3}
    assert sections[2].metadata["table"] == "表格 1"
    assert sections[2].metadata["cell_range"] == "A1:B1"


def test_parse_spreadsheet_and_csv_preserve_cell_ranges(tmp_path) -> None:
    from openpyxl import Workbook

    workbook_path = tmp_path / "roles.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "权限矩阵"
    sheet.append(["角色", "读取", "编辑"])
    sheet.append(["reader", "是", "否"])
    workbook.save(workbook_path)

    spreadsheet_sections = parse_document_sections(workbook_path)
    assert spreadsheet_sections[0].metadata == {
        "sheet": "权限矩阵",
        "row_start": 1,
        "row_end": 1,
        "cell_range": "A1:C1",
    }
    assert spreadsheet_sections[1].metadata["cell_range"] == "A2:C2"

    csv_path = tmp_path / "roles.csv"
    csv_path.write_text("角色,读取,编辑\nreader,是,否\n", encoding="utf-8")
    csv_sections = parse_document_sections(csv_path)
    assert csv_sections[1].metadata["cell_range"] == "A2:C2"
